"""
Générateur principal de templates
"""

import os
from pathlib import Path
from typing import Dict, List, Optional
from loguru import logger

from backend.config import PathConfig
from backend.models.template_config import TemplateConfig
from backend.generator.excel_generator import ExcelTemplateGenerator
from backend.generator.config_generator import ConfigGenerator


class TemplateGenerator:
    """
    Générateur principal de templates livrables.
    """
    
    def __init__(self, template_config: TemplateConfig):
        """
        Initialise le générateur.
        
        Args:
            template_config: Configuration validée du template
        """
        self.config = template_config
        self.template_dir = PathConfig.TEMPLATES / self.config.name
        
    def generate(
        self, 
        ppt_source: Optional[Path] = None,
        excel_source: Optional[Path] = None,
        create_new: bool = False
    ) -> Dict[str, Path]:
        """
        Génère un template complet.
        
        Args:
            ppt_source: Fichier PowerPoint source (ou None pour créer vierge)
            excel_source: Fichier Excel source (ou None pour créer vierge)
            create_new: Si True, crée des fichiers vierges même si sources fournies
            
        Returns:
            Dict des chemins créés
        """
        logger.info(f"Création du template livrable '{self.config.name}' (masters + tables demandées)")

        
        # Création du dossier template
        self.template_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Dossier template créé : {self.template_dir}")
        
        created_files = {}
        
        # 1. Génération config.yaml
        config_path = self._generate_config_yaml()
        created_files['config'] = config_path
        
        # 2. Génération Excel template
        excel_generator = ExcelTemplateGenerator(self.config)
        excel_path = excel_generator.generate(
            source_file=excel_source,
            output_dir=self.template_dir,
            create_new=create_new
        )
        created_files['excel'] = excel_path
        
        # 3. Copie ou création PowerPoint + MISE À JOUR DES LIENS EXCEL
        ppt_path = self._handle_ppt_template(ppt_source, create_new, excel_path)
        created_files['ppt'] = ppt_path
        
        # 4. Génération dossier queries/
        queries_dir = self._generate_queries_directory()
        created_files['queries'] = queries_dir
        
        # 5. Génération README
        readme_path = self._generate_readme()
        created_files['readme'] = readme_path
        
        logger.success(f"Template '{self.config.name}' généré avec succès")
        logger.info(f"Fichiers créés :")
        for key, path in created_files.items():
            logger.info(f"  - {key}: {path}")
        
        return created_files
    
    def _generate_config_yaml(self) -> Path:
        """Génère le fichier config.yaml"""
        config_generator = ConfigGenerator(self.config)
        config_path = self.template_dir / "config.yaml"
        config_generator.generate(config_path)
        return config_path
    
    def _handle_ppt_template(
        self, 
        source: Optional[Path], 
        create_new: bool,
        excel_path: Path
    ) -> Path:
        """
        Gère la création/copie du template PowerPoint.
        CORRECTION : Met à jour les liens Excel vers le nouveau master.xlsx
        """
        ppt_path = self.template_dir / "master.pptx"
        
        if create_new or source is None:
            logger.info("Création d'un PowerPoint vierge")
            from pptx import Presentation
            prs = Presentation()
            prs.slide_width = 9144000  # 25.4 cm
            prs.slide_height = 6858000  # 19.05 cm
            
            # Slide de titre
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = self.config.name
            
            prs.save(str(ppt_path))
        else:
            logger.info(f"Copie du PowerPoint source : {source}")
            import shutil
            shutil.copy2(source, ppt_path)
            
            # CORRECTION : Mettre à jour les liens Excel
            self._update_excel_links_in_ppt(ppt_path, excel_path)
        
        return ppt_path
    
    def _update_excel_links_in_ppt(self, ppt_path: Path, excel_path: Path) -> None:
        """
        Met à jour tous les liens Excel dans le PowerPoint pour pointer vers le nouveau master.xlsx.
        Utilise la même méthode que report_service.py qui fonctionne.
        """
        from backend.core.ppt_handler import powerpoint_app_context
        
        logger.info("Mise à jour des liens Excel dans le PowerPoint...")
        
        excel_path_abs = os.path.abspath(str(excel_path))
        
        with powerpoint_app_context(str(ppt_path), visible=True) as (ppt_app, presentation):
            updated_links = 0
            
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    try:
                        if shape.Type == 3 and hasattr(shape, 'LinkFormat') and shape.LinkFormat:
                            old_source = shape.LinkFormat.SourceFullName
                            shape.LinkFormat.SourceFullName = excel_path_abs
                            updated_links += 1
                            logger.debug(f"Lien mis à jour : {old_source} -> {excel_path_abs}")
                    except:
                        continue
            
            if updated_links > 0:
                logger.success(f"{updated_links} lien(s) Excel mis à jour")
                presentation.Save()
            else:
                logger.warning("Aucun lien Excel trouvé dans le PowerPoint")
    
    def _generate_queries_directory(self) -> Path:
        """Génère le dossier queries/ avec templates SQL"""
        queries_dir = self.template_dir / "queries"
        queries_dir.mkdir(exist_ok=True)
        
        # Génération d'un fichier SQL template par table
        for table in self.config.data_source.required_tables:
            sql_file = queries_dir / f"{table}.sql"
            self._generate_sql_template(sql_file, table)
        
        return queries_dir
    
    def _generate_sql_template(self, sql_file: Path, table_name: str) -> None:
        """Génère un template SQL"""
        # Construire les paramètres pour la requête
        param_placeholders = ", ".join([f"{{{p.name}}}" for p in self.config.parameters])
        
        sql_content = f"""-- Template SQL pour {table_name}
-- Généré automatiquement par KAIVAA Builder

-- Paramètres disponibles :
{chr(10).join([f"-- - {{{p.name}}} : {p.description or p.type}" for p in self.config.parameters])}

SELECT 
    *
FROM {table_name}
WHERE 1=1
{chr(10).join([f"  AND {p.name} = '{{{p.name}}}'" for p in self.config.parameters if p.type == 'string'])}
ORDER BY created_at DESC;
"""
        
        with open(sql_file, 'w', encoding='utf-8') as f:
            f.write(sql_content)
    
    def _generate_readme(self) -> Path:
        """Génère un README pour le template livrable (masters + tables demandées)"""
        readme_path = self.template_dir / "README.md"

        params_list = "\n".join([
            f"- **{p.name}** ({p.type}): {p.description or 'Pas de description'}" +
            (" - Obligatoire" if p.required else "")
            for p in self.config.parameters
        ])

        params_example = "\n".join([
            f'        "{p.name}": "valeur",'
            for p in self.config.parameters
        ])

        readme_lines = [
            f"# {self.config.name}",
            "",
            f"**Version:** {self.config.version}",
            f"**Créé le:** {self.config.created_at.strftime('%Y-%m-%d')}",
        ]

        if self.config.created_by:
            readme_lines.append(f"**Par:** {self.config.created_by}")

        readme_lines.extend([
            "",
            "## Description",
            "",
            self.config.description or "Pas de description",
            "",
            "## Paramètres",
            "",
            params_list,
            "",
            "## Tables demandées",
            "",
            f"- Type: {self.config.data_source.type}",
            f"- Tables demandées (par gabarit): {', '.join(self.config.data_source.required_tables)}",
            "",
            "## Structure des fichiers",
            "",
            "```",
            f"{self.config.name}/",
            "├── config.yaml           # Configuration du template livrable (tables demandées)",
            "├── master.pptx          # Master PPT (facultatif)",
            "├── master.xlsx          # Master Excel (obligatoire)",
            "├── queries/             # Requêtes SQL",
            "│   ├── table1.sql",
            "│   └── table2.sql",
            "└── README.md            # Ce fichier",
            "```",
            "",
            "## Utilisation",
            "",
            "Pour générer un rapport avec ce template :",
            "",
            "```python",
            "from backend.services.report_service import ReportService",
            "",
            "service = ReportService()",
            "result = service.generate_report(",
            f'    template_name="{self.config.name}",',
            "    parameters={",
            params_example,
            "    }",
            ")",
            "```"
        ])

        readme_content = "\n".join(readme_lines)

        with open(readme_path, 'w', encoding='utf-8') as f:
            f.write(readme_content)

        return readme_path
